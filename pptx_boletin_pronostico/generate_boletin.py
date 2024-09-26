import os
import pptx
import datetime
import locale

from pptx.dml.color import RGBColor

periodo = 'SON'
anio = '2024'

# Establecer la configuración regional en español
locale.setlocale(locale.LC_TIME, 'es_ES.UTF-8')
# Diccionario que mapea las iniciales de los meses en español a sus nombres completos
meses = {
    'DEF': 'Enero',
    'EFM': 'Febrero',
    'FMA': 'Marzo',  # 'M' podría ser Marzo o Mayo, se debe manejar adecuadamente
    'MAM': 'Abril',
    'AMJ': 'Mayo',   # 'M' podría ser Marzo o Mayo, se debe manejar adecuadamente
    'MJJ': 'Junio',
    'JJA': 'Julio',  # 'J' podría ser Junio o Julio, se debe manejar adecuadamente
    'JAS': 'Agosto',
    'ASO': 'Septiembre',
    'SON': 'Octubre',
    'OND': 'Noviembre',
    'NDE': 'Diciembre'
}
# Lista de meses para resolver ambigüedades
meses_ordenados = ['Enero', 'Febrero', 'Marzo', 'Abril', 'Mayo', 'Junio', 'Julio', 'Agosto', 'Septiembre', 'Octubre', 'Noviembre', 'Diciembre']

def obtener_periodo(periodo, anio):
    # Inicializar una lista para almacenar los nombres de los meses del periodo
    fecha_actual = datetime.datetime.now()

    fecha_str = '01' + meses[periodo] + anio
    fecha = datetime.datetime.strptime(fecha_str, '01%B%Y')

    mes_anterior = fecha - datetime.timedelta(days=1)
    mes_siguiente = fecha + datetime.timedelta(days=20)

    nombre_mes_anterior = mes_anterior.strftime('%B').capitalize()
    nombre_mes_siguiente = mes_siguiente.strftime('%B').capitalize()

    #Ej Mayo Julio, 2024             fecha actual      mes siguiente
    return [nombre_mes_anterior + ' - ' + nombre_mes_siguiente + ', ' + anio,
            fecha_actual.strftime("%d de %b del %Y"), nombre_mes_siguiente]


dir_base = 'D:/CIIFEN/Boletines_Sudamerica/'
ppt = "D:/CIIFEN/Boletines_Sudamerica/Pronostico_Estacional_Sudamerica.pptx"
ruta = dir_base + anio+'/' + periodo + '/UNION/' + 'PNG' + '/'

periodos_text = obtener_periodo(periodo, anio)[0]
date_today = obtener_periodo(periodo, anio)[1]
prox_publ = 'segunda semana de ' + obtener_periodo(periodo, anio)[2]

prs = pptx.Presentation(ppt)

#Info del slide texto por cambiar
texto_por_reemplazar = [[date_today,0,14,0], # fecha completa 30 abril 2024
                        [periodos_text,0,5,0], # abril junio
                        [periodos_text,1,1,1], # slide 2 encabezado
                            [periodos_text,3,4,1], # slide 4 encabezado
                            [periodos_text,2,4,1], # slide 3 encabezado
                            [periodos_text,4,5,1], # slide 5 encabezado
                            [periodos_text,5,5,1], # slide 6 encabezado
                            [prox_publ,6,6,1] # segunda semana de mayo slide 7
                        ]

for info in texto_por_reemplazar:
    # Acceder al texto que deseas actualizar y guardar su formato

    texto_frame = prs.slides[info[1]].shapes[info[2]].text_frame
    parrafo = texto_frame.paragraphs[info[3]]
    #texto_actual = parrafo.text
    # Guardar el formato original del texto
    fuente = parrafo.runs[0].font
    tamanio = fuente.size
    try:
        color = fuente.color.rgb
    except AttributeError:
        color = RGBColor(255, 255, 255)
    italic = fuente.italic
    bold = fuente.bold
    tipo_letra = fuente.name

    # Actualizar el texto manteniendo el formato
    parrafo.clear()  # Limpiar el párrafo existente
    nuevo_run = parrafo.add_run()
    nuevo_run.text = info[0]
    nuevo_run.font.size = tamanio
    nuevo_run.font.color.rgb = color
    nuevo_run.font.bold = bold
    nuevo_run.font.italic = italic
    nuevo_run.font.name = tipo_letra



#[name png, slide, shape]
datos =  [['Pronostico_SUD_precip2.png', 1, 3],
   ['Pronostico_VEN_COL_EC_PREC.png', 2, 5],
     ['Pronostico_PER_BOL_CH_PREC.png', 3, 5],
    ['Pronostico_CH_PREC.png', 3, 6],
    ['Pronostico_SUD_tmax2.png', 4, 6],
    ['Pronostico_SUD_tmin2.png', 5, 6]]

for info in datos:
    ruta_png = ruta + info[0]
    # #Get file to replace
    img_shape = prs.slides[info[1]].shapes[info[2]]
    # create new image part from new image file
    new_pptx_img = pptx.parts.image.Image.from_file(ruta_png)

    # get part and rId from shape we need to change
    slide_part, rId = img_shape.part, img_shape._element.blip_rId
    image_part = slide_part.related_part(rId)

    # overwrite old blob info with new blob info
    image_part.blob = new_pptx_img._blob


ruta_final = dir_base+'Pronostico_Estacional_Sudamerica_' + periodo + '_' + anio + '.pptx'
prs.save(ruta_final)


# text_frame = prs.slides[0].shapes[4].text_frame
# print(text_frame.text)
# for i in prs.slides[1].shapes:
#     print(i)